"""Extract per-slide / per-page text from ML course (req 3) slides into
scripts/_ml_dump/M<N>.txt. Handles both PDF and PPTX.
"""
from __future__ import annotations
import re, sys
from pathlib import Path
from pypdf import PdfReader
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
SLICES_DIR = ROOT / "Credit-Transfer-20260428" / "3-Machine Learning" / "Slices"
OUT_DIR = Path(__file__).resolve().parent / "_ml_dump"
OUT_DIR.mkdir(exist_ok=True)


def dump_pdf(pdf: Path, out: Path, idx: str, title: str) -> int:
    reader = PdfReader(str(pdf))
    n = len(reader.pages)
    with out.open("w", encoding="utf-8") as fh:
        fh.write(f"# M{idx}  {title}\n")
        fh.write(f"# Source: {pdf.name}    Pages: {n}\n\n")
        for i, page in enumerate(reader.pages, 1):
            try:
                text = page.extract_text() or ""
            except Exception as exc:
                text = f"<extract failed: {exc}>"
            text = re.sub(r"\n{3,}", "\n\n", text.strip())
            fh.write(f"=== p{i} ===\n{text}\n\n")
    return n


def dump_pptx(pptx: Path, out: Path, idx: str, title: str) -> int:
    prs = Presentation(str(pptx))
    slides = list(prs.slides)
    n = len(slides)
    with out.open("w", encoding="utf-8") as fh:
        fh.write(f"# M{idx}  {title}\n")
        fh.write(f"# Source: {pptx.name}    Slides: {n}\n\n")
        for i, slide in enumerate(slides, 1):
            chunks = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        line = "".join(r.text for r in p.runs).strip()
                        if line:
                            chunks.append(line)
                elif shape.shape_type == 19 and hasattr(shape, "text"):
                    if shape.text.strip():
                        chunks.append(shape.text.strip())
            text = "\n".join(chunks)
            text = re.sub(r"\n{3,}", "\n\n", text)
            fh.write(f"=== s{i} ===\n{text}\n\n")
    return n


def main() -> None:
    files = sorted(SLICES_DIR.iterdir(), key=lambda p: int(re.match(r"(\d+)", p.name).group(1)) if re.match(r"\d+", p.name) else 999)
    for f in files:
        m = re.match(r"(\d+)[-_](.+)\.(pdf|pptx)$", f.name)
        if not m:
            print(f"skip {f.name}")
            continue
        idx, title, ext = m.group(1), m.group(2), m.group(3).lower()
        out = OUT_DIR / f"M{idx}.txt"
        try:
            if ext == "pdf":
                n = dump_pdf(f, out, idx, title)
            else:
                n = dump_pptx(f, out, idx, title)
            print(f"M{idx}: {n:3d} -> {out.name}")
        except Exception as exc:
            print(f"M{idx} FAILED: {exc}")


if __name__ == "__main__":
    main()
